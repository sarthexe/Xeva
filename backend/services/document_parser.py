"""
Document Parser Service - Extract text from various file types

Supports:
- PDF files (with OCR fallback for scanned documents)
- Word documents (.docx)
- Excel spreadsheets (.xlsx)
- PowerPoint presentations (.pptx)
- Images (.png, .jpg, .jpeg, .gif, .bmp, .webp) with OCR
- Plain text (.txt, .md, .csv, .json)
"""

import io
import os
import re
import logging
from typing import Tuple, Optional, List, Dict, Any
from pathlib import Path
from datetime import datetime

# Configure logging
logger = logging.getLogger(__name__)

# PDF parsing
try:
    from PyPDF2 import PdfReader
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False
    logger.warning("PyPDF2 not available - PDF parsing disabled")

# Word document parsing
try:
    from docx import Document as DocxDocument
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    logger.warning("python-docx not available - DOCX parsing disabled")

# Excel spreadsheet parsing
try:
    from openpyxl import load_workbook
    XLSX_AVAILABLE = True
except ImportError:
    XLSX_AVAILABLE = False
    logger.warning("openpyxl not available - Excel parsing disabled")

# PowerPoint parsing
try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logger.warning("python-pptx not available - PowerPoint parsing disabled")

# Image OCR
try:
    from PIL import Image
    import pytesseract
    OCR_AVAILABLE = True
except ImportError:
    OCR_AVAILABLE = False
    logger.warning("Pillow/pytesseract not available - Image OCR disabled")

# OpenCV for image preprocessing
try:
    import cv2
    import numpy as np
    CV2_AVAILABLE = True
except ImportError:
    CV2_AVAILABLE = False
    logger.warning("OpenCV not available - Image preprocessing disabled")

# PDF to image for OCR fallback
try:
    from pdf2image import convert_from_bytes
    PDF2IMAGE_AVAILABLE = True
except ImportError:
    PDF2IMAGE_AVAILABLE = False
    logger.warning("pdf2image not available - PDF OCR fallback disabled")


class DocumentParser:
    """
    Universal document parser for extracting text from various file formats.
    
    Features:
    - Multi-format support (PDF, DOCX, XLSX, PPTX, images, text)
    - Rich metadata extraction
    - OCR with image preprocessing for better accuracy
    - Text chunking utilities for RAG systems
    """
    
    # Supported file extensions by category
    SUPPORTED_TEXT = ['.txt', '.md', '.csv', '.json', '.xml', '.html', '.htm']
    SUPPORTED_PDF = ['.pdf']
    SUPPORTED_DOCX = ['.docx']
    SUPPORTED_XLSX = ['.xlsx']
    SUPPORTED_PPTX = ['.pptx']
    SUPPORTED_IMAGES = ['.png', '.jpg', '.jpeg', '.gif', '.bmp', '.webp', '.tiff', '.tif']
    
    @classmethod
    def get_supported_extensions(cls) -> list:
        """Get all supported file extensions."""
        extensions = cls.SUPPORTED_TEXT.copy()
        if PDF_AVAILABLE:
            extensions.extend(cls.SUPPORTED_PDF)
        if DOCX_AVAILABLE:
            extensions.extend(cls.SUPPORTED_DOCX)
        if XLSX_AVAILABLE:
            extensions.extend(cls.SUPPORTED_XLSX)
        if PPTX_AVAILABLE:
            extensions.extend(cls.SUPPORTED_PPTX)
        if OCR_AVAILABLE:
            extensions.extend(cls.SUPPORTED_IMAGES)
        return extensions
    
    @classmethod
    def is_supported(cls, filename: str) -> bool:
        """Check if a file type is supported."""
        ext = cls._get_extension(filename)
        return ext in cls.get_supported_extensions()
    
    @staticmethod
    def _get_extension(filename: str) -> str:
        """Get lowercase file extension."""
        return Path(filename).suffix.lower()
    
    @classmethod
    async def parse(cls, content: bytes, filename: str) -> Tuple[str, dict]:
        """
        Parse document content and extract text.
        
        Args:
            content: Raw file bytes
            filename: Original filename (used to detect file type)
            
        Returns:
            Tuple of (extracted_text, metadata)
        """
        ext = cls._get_extension(filename)
        metadata = {
            "filename": filename,
            "file_type": ext,
            "file_size_bytes": len(content),
            "parser_used": None,
            "parsed_at": datetime.now().isoformat()
        }
        
        try:
            # Plain text files
            if ext in cls.SUPPORTED_TEXT:
                text = cls._parse_text(content)
                metadata["parser_used"] = "text"
                return text.strip(), metadata
            
            # PDF files
            elif ext in cls.SUPPORTED_PDF:
                if not PDF_AVAILABLE:
                    raise ValueError("PDF parsing not available. Install PyPDF2.")
                text, pdf_meta = cls._parse_pdf(content)
                metadata.update(pdf_meta)
                metadata["parser_used"] = "pdf"
                return text.strip(), metadata
            
            # Word documents
            elif ext in cls.SUPPORTED_DOCX:
                if not DOCX_AVAILABLE:
                    raise ValueError("DOCX parsing not available. Install python-docx.")
                text, docx_meta = cls._parse_docx(content)
                metadata.update(docx_meta)
                metadata["parser_used"] = "docx"
                return text.strip(), metadata
            
            # Excel spreadsheets
            elif ext in cls.SUPPORTED_XLSX:
                if not XLSX_AVAILABLE:
                    raise ValueError("Excel parsing not available. Install openpyxl.")
                text, xlsx_meta = cls._parse_xlsx(content)
                metadata.update(xlsx_meta)
                metadata["parser_used"] = "xlsx"
                return text.strip(), metadata
            
            # PowerPoint presentations
            elif ext in cls.SUPPORTED_PPTX:
                if not PPTX_AVAILABLE:
                    raise ValueError("PowerPoint parsing not available. Install python-pptx.")
                text, pptx_meta = cls._parse_pptx(content)
                metadata.update(pptx_meta)
                metadata["parser_used"] = "pptx"
                return text.strip(), metadata
            
            # Images (OCR)
            elif ext in cls.SUPPORTED_IMAGES:
                if not OCR_AVAILABLE:
                    raise ValueError("Image OCR not available. Install Pillow and pytesseract.")
                text, img_meta = cls._parse_image(content)
                metadata.update(img_meta)
                metadata["parser_used"] = "ocr"
                return text.strip(), metadata
            
            else:
                raise ValueError(f"Unsupported file type: {ext}")
                
        except Exception as e:
            logger.error(f"Error parsing {filename}: {str(e)}")
            raise ValueError(f"Error parsing {filename}: {str(e)}")
    
    @staticmethod
    def _parse_text(content: bytes) -> str:
        """Parse plain text files with encoding detection."""
        # Try common encodings
        encodings = ['utf-8', 'utf-16', 'latin-1', 'cp1252', 'ascii']
        
        for encoding in encodings:
            try:
                return content.decode(encoding)
            except (UnicodeDecodeError, LookupError):
                continue
        
        # Fallback: decode with errors replaced
        return content.decode('utf-8', errors='replace')
    
    @staticmethod
    def _parse_pdf(content: bytes) -> Tuple[str, dict]:
        """Extract text and metadata from PDF file."""
        pdf_file = io.BytesIO(content)
        reader = PdfReader(pdf_file)
        
        text_parts = []
        total_pages = len(reader.pages)
        
        for page_num, page in enumerate(reader.pages):
            page_text = page.extract_text()
            if page_text:
                text_parts.append(page_text)
        
        extracted_text = "\n\n".join(text_parts)
        
        # Extract PDF metadata
        pdf_metadata = {}
        if reader.metadata:
            meta = reader.metadata
            if meta.author:
                pdf_metadata["author"] = meta.author
            if meta.title:
                pdf_metadata["title"] = meta.title
            if meta.subject:
                pdf_metadata["subject"] = meta.subject
            if meta.creator:
                pdf_metadata["creator"] = meta.creator
            if meta.creation_date:
                try:
                    pdf_metadata["creation_date"] = str(meta.creation_date)
                except:
                    pass
        
        # If no text extracted, try OCR fallback
        if not extracted_text.strip() and PDF2IMAGE_AVAILABLE and OCR_AVAILABLE:
            try:
                images = convert_from_bytes(content)
                ocr_parts = []
                for img in images:
                    # Apply preprocessing if available
                    if CV2_AVAILABLE:
                        img = DocumentParser._preprocess_image_for_ocr(img)
                    ocr_text = pytesseract.image_to_string(img)
                    if ocr_text.strip():
                        ocr_parts.append(ocr_text)
                extracted_text = "\n\n".join(ocr_parts)
                return extracted_text, {
                    "pages": total_pages,
                    "ocr_used": True,
                    "preprocessing_applied": CV2_AVAILABLE,
                    **pdf_metadata
                }
            except Exception as e:
                logger.warning(f"PDF OCR fallback failed: {e}")
        
        return extracted_text, {
            "pages": total_pages,
            "ocr_used": False,
            **pdf_metadata
        }
    
    @staticmethod
    def _parse_docx(content: bytes) -> Tuple[str, dict]:
        """Extract text and metadata from Word document."""
        docx_file = io.BytesIO(content)
        doc = DocxDocument(docx_file)
        
        text_parts = []
        
        # Extract paragraphs
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)
        
        # Extract tables
        table_count = 0
        for table in doc.tables:
            table_count += 1
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    text_parts.append(row_text)
        
        # Extract document properties
        doc_metadata = {
            "paragraphs": len(doc.paragraphs),
            "tables": table_count
        }
        
        core_props = doc.core_properties
        if core_props:
            if core_props.author:
                doc_metadata["author"] = core_props.author
            if core_props.title:
                doc_metadata["title"] = core_props.title
            if core_props.subject:
                doc_metadata["subject"] = core_props.subject
            if core_props.created:
                doc_metadata["created"] = str(core_props.created)
            if core_props.modified:
                doc_metadata["modified"] = str(core_props.modified)
        
        return "\n\n".join(text_parts), doc_metadata
    
    @staticmethod
    def _parse_xlsx(content: bytes) -> Tuple[str, dict]:
        """Extract text from Excel spreadsheet."""
        xlsx_file = io.BytesIO(content)
        workbook = load_workbook(xlsx_file, data_only=True)
        
        text_parts = []
        total_cells = 0
        sheet_names = workbook.sheetnames
        
        for sheet_name in sheet_names:
            sheet = workbook[sheet_name]
            text_parts.append(f"=== Sheet: {sheet_name} ===")
            
            for row in sheet.iter_rows(values_only=True):
                # Filter out None values and convert to strings
                row_values = [str(cell) if cell is not None else "" for cell in row]
                # Skip completely empty rows
                if any(v.strip() for v in row_values):
                    row_text = " | ".join(row_values)
                    text_parts.append(row_text)
                    total_cells += len([v for v in row_values if v.strip()])
            
            text_parts.append("")  # Add blank line between sheets
        
        return "\n".join(text_parts), {
            "sheets": len(sheet_names),
            "sheet_names": sheet_names,
            "total_cells": total_cells
        }
    
    @staticmethod
    def _parse_pptx(content: bytes) -> Tuple[str, dict]:
        """Extract text from PowerPoint presentation."""
        pptx_file = io.BytesIO(content)
        presentation = Presentation(pptx_file)
        
        text_parts = []
        slide_count = 0
        
        for slide_num, slide in enumerate(presentation.slides, 1):
            slide_count += 1
            slide_text = []
            slide_text.append(f"=== Slide {slide_num} ===")
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
                
                # Extract text from tables in slides
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = " | ".join(
                            cell.text.strip() for cell in row.cells if cell.text.strip()
                        )
                        if row_text:
                            slide_text.append(row_text)
            
            if len(slide_text) > 1:  # More than just the header
                text_parts.extend(slide_text)
                text_parts.append("")  # Blank line between slides
        
        # Extract presentation properties
        pptx_metadata = {
            "slides": slide_count
        }
        
        core_props = presentation.core_properties
        if core_props:
            if core_props.author:
                pptx_metadata["author"] = core_props.author
            if core_props.title:
                pptx_metadata["title"] = core_props.title
            if core_props.subject:
                pptx_metadata["subject"] = core_props.subject
            if core_props.created:
                pptx_metadata["created"] = str(core_props.created)
        
        return "\n".join(text_parts), pptx_metadata
    
    @staticmethod
    def _preprocess_image_for_ocr(image) -> Image.Image:
        """
        Preprocess image for better OCR accuracy.
        
        Applies:
        - Grayscale conversion
        - Contrast enhancement (adaptive thresholding)
        - Noise reduction
        """
        # Convert PIL Image to numpy array
        if isinstance(image, Image.Image):
            img_array = np.array(image)
        else:
            img_array = image
        
        # Convert to grayscale if needed
        if len(img_array.shape) == 3:
            gray = cv2.cvtColor(img_array, cv2.COLOR_RGB2GRAY)
        else:
            gray = img_array
        
        # Apply Gaussian blur to reduce noise
        blurred = cv2.GaussianBlur(gray, (3, 3), 0)
        
        # Apply adaptive thresholding for better contrast
        thresh = cv2.adaptiveThreshold(
            blurred, 255,
            cv2.ADAPTIVE_THRESH_GAUSSIAN_C,
            cv2.THRESH_BINARY,
            11, 2
        )
        
        # Apply morphological operations to clean up
        kernel = np.ones((1, 1), np.uint8)
        cleaned = cv2.morphologyEx(thresh, cv2.MORPH_CLOSE, kernel)
        
        # Convert back to PIL Image
        return Image.fromarray(cleaned)
    
    @staticmethod
    def _parse_image(content: bytes) -> Tuple[str, dict]:
        """Extract text from image using OCR with preprocessing."""
        image = Image.open(io.BytesIO(content))
        
        # Get image info
        width, height = image.size
        format_type = image.format
        
        # Apply preprocessing if OpenCV is available
        preprocessing_applied = False
        if CV2_AVAILABLE:
            try:
                processed_image = DocumentParser._preprocess_image_for_ocr(image)
                preprocessing_applied = True
            except Exception as e:
                logger.warning(f"Image preprocessing failed, using original: {e}")
                processed_image = image
        else:
            processed_image = image
        
        # Perform OCR
        try:
            text = pytesseract.image_to_string(processed_image)
        except Exception as e:
            if "tesseract is not installed" in str(e).lower():
                raise ValueError("Tesseract OCR is not installed on the server. Please install it to process images.")
            raise ValueError(f"OCR failed: {str(e)}")
        
        return text, {
            "image_width": width,
            "image_height": height,
            "image_format": format_type,
            "preprocessing_applied": preprocessing_applied
        }
    
    @classmethod
    def chunk_text(
        cls,
        text: str,
        chunk_size: int = 1000,
        overlap: int = 200,
        respect_sentences: bool = True
    ) -> List[Dict[str, Any]]:
        """
        Split text into overlapping chunks for RAG systems.
        
        Args:
            text: The text to chunk
            chunk_size: Target size of each chunk in characters
            overlap: Number of characters to overlap between chunks
            respect_sentences: If True, try not to break mid-sentence
            
        Returns:
            List of dicts with 'text', 'start_pos', 'end_pos', 'chunk_index'
        """
        if not text or not text.strip():
            return []
        
        chunks = []
        text = text.strip()
        
        if len(text) <= chunk_size:
            return [{
                "text": text,
                "start_pos": 0,
                "end_pos": len(text),
                "chunk_index": 0
            }]
        
        # Sentence-ending pattern
        sentence_end = re.compile(r'[.!?]\s+')
        
        start = 0
        chunk_index = 0
        
        while start < len(text):
            end = min(start + chunk_size, len(text))
            
            # If we're not at the end and should respect sentences
            if end < len(text) and respect_sentences:
                # Look for sentence end near the chunk boundary
                search_start = max(start + chunk_size - 100, start)
                search_region = text[search_start:end + 50]
                
                matches = list(sentence_end.finditer(search_region))
                if matches:
                    # Use the last sentence end in the region
                    last_match = matches[-1]
                    end = search_start + last_match.end()
            
            chunk_text = text[start:end].strip()
            
            if chunk_text:
                chunks.append({
                    "text": chunk_text,
                    "start_pos": start,
                    "end_pos": end,
                    "chunk_index": chunk_index
                })
                chunk_index += 1
            
            # Move start with overlap
            start = end - overlap if end < len(text) else end
            
            # Prevent infinite loop
            if start >= end:
                break
        
        return chunks


# Singleton instance
document_parser = DocumentParser()
